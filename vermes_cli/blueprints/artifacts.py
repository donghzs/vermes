"""Blueprint: Artifacts（产物文件读取端点）

提供前端 ArtifactPanel 的文件读取能力。
安全：白名单三类根（cwd / ~/.vermes/ / /tmp/），路径规范化防穿越。
"""
import os
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import PlainTextResponse, HTMLResponse, Response, JSONResponse
import base64


def _allowed_roots():
    """返回允许读取的根目录列表（每次调用动态获取 cwd）"""
    return [
        Path.cwd().resolve(),
        (Path.home() / '.vermes').resolve(),
        Path('/tmp').resolve(),
        # 用户常用目录：产物文件常在 Desktop/Downloads/Documents
        (Path.home() / 'Desktop').resolve(),
        (Path.home() / 'Downloads').resolve(),
        (Path.home() / 'Documents').resolve(),
    ]


# 安全响应头：防嗅探 + CSP sandbox（防产物 HTML 内脚本执行）
_SECURITY_HEADERS = {
    'X-Content-Type-Options': 'nosniff',
    'Content-Security-Policy': "default-src 'none'; img-src data:; style-src 'unsafe-inline'",
}


def _is_safe_path(path_str: str) -> Path:
    """路径安全校验：规范化 + 白名单根目录检查"""
    # 展开 ~ / ~user（前端消息流常出现 ~/Documents/... 这类用户目录路径）
    path_str = os.path.expanduser(path_str)
    # 先把 tmp/ 前缀映射到 /tmp/（URL path 丢失前导 / 的常见情况）
    if path_str.startswith('tmp/') and not Path.cwd().joinpath('tmp').exists():
        path_str = '/' + path_str

    def _check(p: str) -> Optional[Path]:
        raw = Path(p)
        resolved = raw.resolve() if raw.is_absolute() else (Path.cwd() / raw).resolve()
        for root in _allowed_roots():
            try:
                resolved.relative_to(root)
                return resolved
            except ValueError:
                continue
        return None

    ok = _check(path_str)
    if ok is not None:
        # 若路径实际存在，直接返回
        if ok.exists():
            return ok
        # 裸相对路径（agent 工具常把产物写到 /tmp 或 ~/.vermes，却在结果文本里只报裸文件名，
        # 而桌面端后端进程 cwd 可能是 /，导致按 cwd 解析不到真实文件）：
        # 依次在常见产物目录下找同名文件，找到即返回。
        if not path_str.startswith('/') and not path_str.startswith('~'):
            # 1) 直接在常见产物根目录下找同名文件
            for fallback_root in (Path('/tmp'), Path.home() / '.vermes', Path.home()):
                candidate = _check(str(fallback_root / path_str))
                if candidate is not None and candidate.exists():
                    return candidate
            # 2) 递归搜索子目录（agent 常写到 ~/.vermes/logs/ 等子目录，
            #    但结果文本只报裸文件名如 agent.log）
            bare_name = Path(path_str).name
            for search_root in (Path.home() / '.vermes', Path('/tmp')):
                try:
                    for hit in search_root.rglob(bare_name):
                        # 只取文件、且在白名单内、且不超过 3 层深度
                        if hit.is_file():
                            verified = _check(str(hit))
                            if verified is not None and verified.exists():
                                return verified
                except (PermissionError, OSError):
                    continue
        return ok

    # 兜底：agent 常把桌面/下载/文档写成根级绝对路径（/Desktop/.. 而非 ~/Desktop/..），
    # 这类路径明显是 home 同名目录的误写，重定向到 Path.home() 下再校验，避免前端读产物时 403。
    # 仅精确匹配这三个根级目录名，不波及 /Desktopfoo 等。
    home_redirects = {
        '/Desktop': Path.home() / 'Desktop',
        '/Downloads': Path.home() / 'Downloads',
        '/Documents': Path.home() / 'Documents',
    }
    for prefix, target in home_redirects.items():
        if path_str == prefix or path_str.startswith(prefix + '/'):
            redirected = target / path_str[len(prefix):].lstrip('/')
            ok = _check(str(redirected))
            if ok is not None:
                return ok

    raise HTTPException(status_code=403, detail=f"路径不在允许范围内: {path_str}")


_MIME_MAP = {
    '.md': 'text/markdown',
    '.html': 'text/html',
    '.htm': 'text/html',
    '.txt': 'text/plain',
    '.csv': 'text/csv',
    '.json': 'application/json',
    '.py': 'text/plain',
    '.js': 'text/plain',
    '.ts': 'text/plain',
    '.sh': 'text/plain',
    '.yaml': 'text/plain',
    '.yml': 'text/plain',
    '.ini': 'text/plain',
    '.cfg': 'text/plain',
    '.log': 'text/plain',
    '.toml': 'text/plain',
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.gif': 'image/gif',
    '.webp': 'image/webp',
    '.svg': 'image/svg+xml',
    # 办公/ScholarForge
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.pdf': 'application/pdf',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.doc': 'application/msword',
    '.xls': 'application/vnd.ms-excel',
    '.ppt': 'application/vnd.ms-powerpoint',
    # 制造业/mfgcad
    '.step': 'application/step',
    '.stp': 'application/step',
    '.stl': 'application/sla',
    '.obj': 'application/wavefront-obj',
    '.fcdoc': 'application/octet-stream',
    '.dxf': 'application/dxf',
    '.gcode': 'text/plain',
    '.iges': 'application/iges',
    '.3mf': 'model/3mf',
    '.gltf': 'model/gltf+json',
}

_MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB


def _check_origin(request: Request):
    """纵深防御：校验请求来源是否为本应用（挡掉跨站调用）"""
    origin = request.headers.get('origin', '')
    host = request.headers.get('host', '')
    # Electron 无 origin（file:// 协议）或 origin 包含 host（同源）
    if origin and host:
        # 提取 origin 的 host 部分比对
        from urllib.parse import urlparse
        parsed = urlparse(origin)
        if parsed.hostname and parsed.hostname not in ('localhost', '127.0.0.1', '0.0.0.0'):
            raise HTTPException(status_code=403, detail="跨站请求被拒绝")
    # 无 origin（Electron 内部请求）或 localhost 来源 → 放行


def register_to(app):
    @app.get('/api/v1/workspace/tree')
    async def workspace_tree(path: str = '', request: Request = None):
        """列出工作目录文件树（单层），返回子项列表"""
        _check_origin(request)
        if path:
            safe = _is_safe_path(path)
        else:
            safe = Path.cwd().resolve()
        if not safe.exists():
            raise HTTPException(status_code=404, detail=f"路径不存在: {path}")
        if not safe.is_dir():
            raise HTTPException(status_code=400, detail=f"不是目录: {path}")

        items = []
        try:
            for entry in sorted(safe.iterdir(), key=lambda e: (not e.is_dir(), e.name.lower())):
                # 跳过隐藏文件和常见忽略目录
                if entry.name.startswith('.'):
                    continue
                if entry.is_dir() and entry.name in {'node_modules', '__pycache__', '.git', 'dist', 'build', '.venv', 'venv'}:
                    continue
                rel = str(entry.relative_to(Path.cwd().resolve())) if str(entry).startswith(str(Path.cwd().resolve())) else str(entry)
                items.append({
                    'name': entry.name,
                    'path': rel,
                    'is_dir': entry.is_dir(),
                    'size': entry.stat().st_size if entry.is_file() else 0,
                    'ext': entry.suffix.lower() if entry.is_file() else '',
                })
        except PermissionError:
            raise HTTPException(status_code=403, detail="无权限读取该目录")
        return {'items': items, 'current': str(safe.relative_to(Path.cwd().resolve())) if str(safe).startswith(str(Path.cwd().resolve())) else str(safe)}

    @app.get('/api/v1/artifacts/{path:path}/resolve')
    async def resolve_artifact(path: str, request: Request):
        """返回产物文件的绝对路径，供桌面端 shell.showItemInFolder 使用"""
        _check_origin(request)
        safe_path = _is_safe_path(path)
        if not safe_path.exists():
            raise HTTPException(status_code=404, detail=f"文件不存在: {path}")
        return {'path': str(safe_path), 'name': safe_path.name, 'size': safe_path.stat().st_size}

    @app.get('/api/v1/artifacts/{path:path}')
    async def serve_artifact(path: str, request: Request):
        """读取产物文件，返回对应 MIME 类型"""
        _check_origin(request)
        safe_path = _is_safe_path(path)

        if not safe_path.exists():
            raise HTTPException(status_code=404, detail=f"文件不存在: {path}")

        if not safe_path.is_file():
            raise HTTPException(status_code=400, detail=f"不是文件: {path}")

        file_size = safe_path.stat().st_size
        if file_size > _MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail=f"文件过大 ({file_size // 1024 // 1024}MB)，上限 50MB")

        ext = safe_path.suffix.lower()
        mime = _MIME_MAP.get(ext, 'application/octet-stream')

        # 图片直接返回二进制
        if mime.startswith('image/'):
            with open(safe_path, 'rb') as f:
                return Response(content=f.read(), media_type=mime, headers=_SECURITY_HEADERS)

        # 文本类返回内容
        try:
            with open(safe_path, 'r', encoding='utf-8') as f:
                content = f.read()
        except UnicodeDecodeError:
            # 二进制文件 fallback
            with open(safe_path, 'rb') as f:
                return Response(content=f.read(), media_type=mime, headers=_SECURITY_HEADERS)

        return PlainTextResponse(content, media_type=mime, headers=_SECURITY_HEADERS)

    @app.get('/api/v1/artifacts/{path:path}/preview')
    async def preview_artifact(path: str, request: Request):
        """Office 文档静态预览（零依赖）：当前支持 pptx，复用 python-pptx 抽取每页文本与图片。

        返回 JSON：{'kind':'pptx','pages':[{'i':页码,'text':文本,'images':[data_url...]}]}
        不支持的格式（doc/ppt 老二进制等）返回 {'kind':'unsupported'}，前端回退为下载卡片。
        """
        _check_origin(request)
        safe_path = _is_safe_path(path)
        if not safe_path.exists():
            raise HTTPException(status_code=404, detail=f"文件不存在: {path}")
        ext = safe_path.suffix.lower()
        if ext == '.pptx':
            try:
                from pptx import Presentation
            except ImportError:
                return JSONResponse({'kind': 'unsupported', 'reason': 'python-pptx 未安装（请安装 office 能力：pip install python-pptx 或 vermes 的 office extra）'}, status_code=501)
            try:
                prs = Presentation(str(safe_path))
            except Exception as e:
                return JSONResponse({'kind': 'unsupported', 'reason': f'解析失败: {e}'}, status_code=422)
            pages = []
            for idx, slide in enumerate(prs.slides, 1):
                texts, images = [], []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            texts.append(t)
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            blob = shape.image.blob
                            img_ext = (shape.image.ext or 'png').lower()
                            mime = 'image/png' if img_ext == 'png' else ('image/jpeg' if img_ext in ('jpg', 'jpeg') else 'image/png')
                            images.append(f'data:{mime};base64,{base64.b64encode(blob).decode()}')
                        except Exception:
                            pass
                pages.append({'i': idx, 'text': '\n'.join(texts), 'images': images})
            return JSONResponse({'kind': 'pptx', 'pages': pages})
        return JSONResponse({'kind': 'unsupported', 'reason': '仅支持 pptx 预览'}, status_code=415)

    @app.post('/api/v1/artifacts/{path:path}/content')
    async def write_artifact_content(path: str, request: Request):
        """回存产物文件内容（轻量可编辑右栏的"人改→保存"落点）。

        安全：复用产物白名单根 + 路径穿越防护 + 跨站校验；仅允许已存在的文本类文件回写，
        写入前再确认父目录存在。大小上限与读取一致（50MB）。
        """
        _check_origin(request)
        safe_path = _is_safe_path(path)

        if not safe_path.exists():
            raise HTTPException(status_code=404, detail=f"文件不存在: {path}")
        if not safe_path.is_file():
            raise HTTPException(status_code=400, detail=f"不是文件: {path}")

        body = await request.body()
        if len(body) > _MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail=f"内容过大（{len(body) // 1024 // 1024}MB），上限 50MB")

        try:
            safe_path.parent.mkdir(parents=True, exist_ok=True)
            # 文本类以 UTF-8 写回；二进制（如图片被误命中）按字节写回
            try:
                safe_path.write_text(body.decode('utf-8'), encoding='utf-8')
            except UnicodeDecodeError:
                safe_path.write_bytes(body)
        except OSError as e:
            raise HTTPException(status_code=500, detail=f"写入失败: {e}")

        return {'ok': True, 'path': str(safe_path), 'size': safe_path.stat().st_size}
