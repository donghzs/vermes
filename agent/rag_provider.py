#!/usr/bin/env python3
"""
RAG Provider — Agent 记忆型 RAG (Retrieval-Augmented Generation)

Uses SQLite + FTS5 for lightweight document indexing and retrieval.
No vector database required — FTS5 trigram tokenizer handles Chinese + English.

Documents are chunked (500 chars, 100 char overlap) and indexed into FTS5.
Prefetch queries the top-K chunks matching the user's message and injects
them into the system prompt as [知识库上下文].
"""

import json
import logging
import os
import re
import sqlite3
import threading
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from agent.memory_provider import MemoryProvider
from hermes_constants import get_hermes_home

logger = logging.getLogger(__name__)

_db_lock = threading.Lock()
_conn_cache: Dict[str, sqlite3.Connection] = {}

# ── Vector backend config (A-1) ──────────────────────────────────────
# Default: "fts5" (zero-dependency, always available).
# Set to "sqlite-vec" to enable vector search with sqlite-vec extension.
# Fail-open: if vector backend fails to load, falls back to FTS5 silently.
_VEC_BACKEND = os.environ.get("VERMES_RAG_BACKEND", "fts5").lower()
_vec_available = False  # set True if vec0.dylib loads successfully
_vec_dylib_path: Optional[str] = None


def _try_init_vec() -> bool:
    """Try to load sqlite-vec extension. Returns True if available.
    Fail-open: any error → return False, FTS5 remains the default."""
    global _vec_available, _vec_dylib_path
    if _vec_available:
        return True
    try:
        import sqlite_vec
        _vec_dylib_path = os.path.join(
            os.path.dirname(sqlite_vec.__file__), "vec0.dylib"
        )
        if not os.path.exists(_vec_dylib_path):
            # On non-macOS, the filename might differ (vec0.so on Linux)
            _vec_dylib_path = os.path.join(
                os.path.dirname(sqlite_vec.__file__), "vec0.so"
            )
        if not os.path.exists(_vec_dylib_path):
            logger.debug("sqlite-vec dylib not found, vector backend disabled")
            return False
        _vec_available = True
        logger.info("sqlite-vec vector backend available: %s", _vec_dylib_path)
        return True
    except ImportError:
        logger.debug("sqlite-vec not installed, vector backend disabled")
        return False
    except Exception as e:
        logger.debug("sqlite-vec init failed: %s", e)
        return False


# Try once at import time (fail-open, no crash)
if _VEC_BACKEND == "sqlite-vec":
    _try_init_vec()


def _get_rag_db() -> Path:
    """Get the RAG database path."""
    return get_hermes_home() / "rag" / "documents.db"


def _get_conn(db_path: str) -> sqlite3.Connection:
    """Return a thread-safe cached connection with WAL + busy_timeout.
    If vector backend is enabled, loads sqlite-vec extension."""
    key = str(db_path)
    with _db_lock:
        if key in _conn_cache:
            try:
                _conn_cache[key].execute("SELECT 1")
                return _conn_cache[key]
            except sqlite3.ProgrammingError:
                pass
        conn = sqlite3.connect(key, check_same_thread=False)
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("PRAGMA synchronous=NORMAL")
        conn.execute("PRAGMA busy_timeout=5000")
        # Load sqlite-vec extension if enabled (A-1)
        if _VEC_BACKEND == "sqlite-vec" and _vec_available and _vec_dylib_path:
            try:
                conn.enable_load_extension(True)
                conn.load_extension(_vec_dylib_path)
                conn.enable_load_extension(False)
            except Exception as e:
                logger.debug("Failed to load sqlite-vec on %s: %s", db_path, e)
                # Fail-open: continue with FTS5 only
        _conn_cache[key] = conn
        return conn


def _init_db(db_path: Path) -> None:
    """Initialize RAG database with FTS5."""
    db_path.parent.mkdir(parents=True, exist_ok=True)
    conn = _get_conn(str(db_path))
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            path TEXT NOT NULL,
            filename TEXT NOT NULL,
            file_type TEXT,
            file_size INTEGER,
            ingested_at TEXT NOT NULL,
            chunk_count INTEGER DEFAULT 0
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS chunks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            doc_id INTEGER NOT NULL,
            chunk_index INTEGER NOT NULL,
            content TEXT NOT NULL,
            char_count INTEGER DEFAULT 0,
            FOREIGN KEY (doc_id) REFERENCES documents(id) ON DELETE CASCADE
        )
    """)
    # FTS5 with trigram tokenizer for CJK + English support
    try:
        c.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS chunks_fts
            USING fts5(content, tokenize='trigram')
        """)
    except sqlite3.OperationalError:
        # FTS5 trigram not available — fall back to unicode61
        logger.warning("FTS5 trigram unavailable, falling back to unicode61")
        c.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS chunks_fts
            USING fts5(content)
        """)
    # Vector index (A-1): only create when sqlite-vec is loaded
    # Default FTS5 path is unaffected — this is additive.
    if _VEC_BACKEND == "sqlite-vec" and _vec_available:
        try:
            c.execute("""
                CREATE VIRTUAL TABLE IF NOT EXISTS chunks_vec
                USING vec0(embedding float[1536])
            """)
            logger.info("Vector index (chunks_vec) initialized with sqlite-vec")
        except Exception as e:
            logger.warning("Failed to create chunks_vec table: %s — vector search disabled", e)
    conn.commit()


def _chunk_text(text: str, chunk_size: int = 0, overlap: int = 0, file_ext: str = '') -> List[str]:
    """Split text into overlapping chunks.

    Auto-selects chunk size based on document type:
    - Academic papers (.pdf, .docx, .tex): 1200 chars, 200 overlap
    - Other documents: 500 chars, 100 overlap
    """
    if chunk_size == 0:
        academic_exts = {'.pdf', '.docx', '.tex', '.rtf'}
        if file_ext.lower() in academic_exts:
            chunk_size, overlap = 1200, 200
        else:
            chunk_size, overlap = 500, 100
    if len(text) <= chunk_size:
        return [text] if text.strip() else []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk)
        start = end - overlap
    return chunks


def index_memory_text(target: str, content: str) -> None:
    """Index curated ``memory``-tool content into the RAG FTS5 store.

    The ``memory`` tool (MemoryStore) persists to MEMORY.md / USER.md — a
    *separate* store from the RAG knowledge base that ``memory_search``
    queries. Previously, writing via ``memory`` and then searching via
    ``memory_search`` returned nothing because the two stores never talked.
    This bridges them: every ``add`` / ``replace`` re-indexes the full target
    content here, so ``memory_search`` can find curated memory.

    Memory stores are small and bounded (~8KB), so a full re-index per write is
    cheap and naturally de-duplicates replaces (old rows are deleted first).

    Fail-open: any error is logged and swallowed so memory writes never break.
    """
    try:
        if not content or not content.strip():
            return
        db_path = _get_rag_db()
        _init_db(db_path)
        conn = _get_conn(str(db_path))
        c = conn.cursor()
        doc_name = f"memory:{target}"
        # Remove any previous index for this target (incl. the FTS shadow rows,
        # which are NOT FK-cascaded), then re-index fresh.
        c.execute("SELECT id FROM documents WHERE filename = ?", (doc_name,))
        for (old_doc_id,) in c.fetchall():
            c.execute(
                "DELETE FROM chunks_fts WHERE rowid IN "
                "(SELECT id FROM chunks WHERE doc_id = ?)",
                (old_doc_id,),
            )
            c.execute("DELETE FROM chunks WHERE doc_id = ?", (old_doc_id,))
        c.execute(
            "INSERT INTO documents (path, filename, file_type, file_size, ingested_at, chunk_count) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (doc_name, doc_name, "memory", len(content), datetime.now().isoformat(), 0),
        )
        doc_id = c.lastrowid
        chunks = _chunk_text(content, file_ext=".md")
        for i, ch in enumerate(chunks):
            c.execute(
                "INSERT INTO chunks (doc_id, chunk_index, content, char_count) VALUES (?, ?, ?, ?)",
                (doc_id, i, ch, len(ch)),
            )
            chunk_id = c.lastrowid
            c.execute("INSERT INTO chunks_fts (rowid, content) VALUES (?, ?)", (chunk_id, ch))
        c.execute("UPDATE documents SET chunk_count = ? WHERE id = ?", (len(chunks), doc_id))
        conn.commit()
    except Exception:
        logger.debug("index_memory_text failed (non-fatal)", exc_info=True)


def _extract_text(file_path: str) -> str:
    """Extract text from a file. Supports txt/md/py/js/json + PDF/DOCX/XLSX."""
    ext = Path(file_path).suffix.lower()
    text_encodings = {'.txt', '.md', '.py', '.js', '.ts', '.json', '.yaml', '.yml',
                      '.html', '.css', '.xml', '.csv', '.tsv', '.sh', '.sql', '.log',
                      '.rtf', '.org', '.rst'}
    if ext in text_encodings or ext == '':
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        except Exception as e:
            logger.error("Failed to read %s: %s", file_path, e)
            return ""
    # Binary formats — use bytes extractor
    try:
        with open(file_path, 'rb') as f:
            raw = f.read()
        return _extract_text_from_bytes(raw, ext)
    except RuntimeError:
        # Explicit, user-facing errors (e.g. scanned PDF without OCR engine)
        # must reach the caller instead of being swallowed as empty text.
        raise
    except Exception as e:
        logger.error("Failed to read binary %s: %s", file_path, e)
        return ""


def _extract_text_from_bytes(raw: bytes, ext: str) -> str:
    """Extract text from binary file content (PDF/DOCX/XLSX/PPTX).

    Args:
        raw: Raw file bytes.
        ext: File extension including dot, e.g. '.pdf'.

    Returns:
        Extracted plain text.
    """
    if ext == '.pdf':
        return _extract_pdf(raw)
    elif ext == '.docx':
        return _extract_docx(raw)
    elif ext == '.xlsx':
        return _extract_xlsx(raw)
    elif ext == '.pptx':
        return _extract_pptx(raw)
    else:
        # Try as text with fallback encoding
        for enc in ('utf-8', 'gbk', 'latin-1'):
            try:
                return raw.decode(enc)
            except Exception:
                continue
        return raw.decode('utf-8', errors='replace')


def _extract_pdf(raw: bytes) -> str:
    """Extract text from PDF bytes using PyMuPDF (fitz), with OCR fallback.

    Returns extracted text. Raises RuntimeError with an explicit Chinese
    message when the PDF appears to be a scanned/image PDF and no OCR engine
    is available to recover its text.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error("PyMuPDF (fitz) not installed — cannot parse PDF")
        return ""

    doc = fitz.open(stream=raw, filetype='pdf')
    try:
        parts = []
        for page in doc:
            parts.append(page.get_text())
        text = '\n\n'.join(parts)
        if text.strip():
            return text
        # Empty text layer → likely a scanned/image-only PDF. Try OCR.
        logger.warning("PDF text layer empty (possibly scanned PDF): %d bytes", len(raw))
        return _extract_pdf_ocr(doc)
    finally:
        doc.close()


def _extract_pdf_ocr(doc) -> str:
    """Best-effort OCR over a PyMuPDF document. Returns text or raises."""
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        raise RuntimeError(
            "PDF 为扫描件/图片型，未检测到 OCR 引擎（需安装 pytesseract 与 "
            "tesseract-ocr 及中文语言包 chi_sim）。请安装 OCR 引擎后重试，或将 "
            "PDF 转为包含可复制文字的版本。"
        )
    try:
        pages_text = []
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            pages_text.append(pytesseract.image_to_string(img, lang="chi_sim+eng"))
        ocr_text = '\n\n'.join(pages_text)
        if ocr_text.strip():
            return ocr_text
        raise RuntimeError(
            "PDF 经 OCR 处理仍未提取到文字，可能是空白页或无法识别的内容。"
        )
    except RuntimeError:
        raise
    except Exception as exc:
        raise RuntimeError(f"PDF OCR 处理失败：{exc}")


def _extract_docx(raw: bytes) -> str:
    """Extract text from DOCX bytes using python-docx."""
    try:
        import io
        from docx import Document
        doc = Document(io.BytesIO(raw))
        parts = []
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Extract table text
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(' | '.join(cells))
        return '\n\n'.join(parts)
    except ImportError:
        logger.error("python-docx not installed — cannot parse DOCX")
        return ""
    except Exception as e:
        logger.error("DOCX extraction error: %s", e)
        return ""


def _extract_xlsx(raw: bytes) -> str:
    """Extract text from XLSX bytes using openpyxl."""
    try:
        import io
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f'## {ws.title}')
            for row in ws.iter_rows(max_row=500, values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(' | '.join(cells))
        wb.close()
        return '\n\n'.join(parts)
    except ImportError:
        logger.error("openpyxl not installed — cannot parse XLSX")
        return ""
    except Exception as e:
        logger.error("XLSX extraction error: %s", e)
        return ""


def _extract_pptx(raw: bytes) -> str:
    """Extract text from PPTX bytes using python-pptx."""
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f'## Slide {i}')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
        return '\n\n'.join(parts)
    except ImportError:
        logger.error("python-pptx not installed — cannot parse PPTX")
        return ""
    except Exception as e:
        logger.error("PPTX extraction error: %s", e)
        return ""


class RAGProvider(MemoryProvider):
    """Agent 记忆型 RAG — FTS5 document indexing + retrieval."""

    def __init__(self):
        self._db_path: Optional[Path] = None
        self._session_id: str = ""
        self._initialized = False

    @property
    def name(self) -> str:
        return "rag"

    def is_available(self) -> bool:
        return True  # Always available — just needs SQLite

    def initialize(self, session_id: str, **kwargs) -> None:
        self._session_id = session_id
        self._db_path = _get_rag_db()
        _init_db(self._db_path)
        self._initialized = True
        logger.info("RAGProvider initialized (db=%s)", self._db_path)

    def system_prompt_block(self) -> str:
        if not self._initialized:
            return ""
        # Count documents
        try:
            conn = _get_conn(str(self._db_path))
            c = conn.cursor()
            c.execute("SELECT COUNT(*) FROM documents")
            doc_count = c.fetchone()[0]
            if doc_count == 0:
                return ""
            return f"[知识库] 已索引 {doc_count} 个文档，可用 memory_search 工具检索。"
        except Exception:
            return ""

    def prefetch(self, query: str, *, session_id: str = "") -> str:
        """Retrieve top-K chunks matching the query via FTS5."""
        if not self._initialized or not query.strip():
            return ""
        try:
            conn = _get_conn(str(self._db_path))
            c = conn.cursor()
            # Build FTS5 query — escape special chars
            safe_query = re.sub(r'[^\w\u4e00-\u9fff\s]', ' ', query).strip()
            if not safe_query:
                return ""
            # FTS5 trigram tokenizer requires ≥3 char substrings.
            # For CJK text, generate 3-char trigrams; for ASCII, use whole words.
            terms = []
            for word in safe_query.split():
                cjk_chars = [ch for ch in word if '\u4e00' <= ch <= '\u9fff']
                ascii_part = ''.join(ch for ch in word if ch not in cjk_chars)
                if ascii_part and len(ascii_part) >= 3:
                    terms.append(ascii_part)
                if len(cjk_chars) >= 3:
                    # Generate 3-char trigrams (sliding window)
                    for i in range(len(cjk_chars) - 2):
                        terms.append(cjk_chars[i] + cjk_chars[i+1] + cjk_chars[i+2])
                # CJK <3 chars: skip (trigram can't match),
                # but if combined with ASCII it might still match
            if not terms:
                return ""
            fts_query = " OR ".join(f'"{t}"' for t in terms[:8])
            c.execute("""
                SELECT chunks.content, documents.filename, chunks.chunk_index
                FROM chunks_fts
                JOIN chunks ON chunks.id = chunks_fts.rowid
                JOIN documents ON documents.id = chunks.doc_id
                WHERE chunks_fts MATCH ?
                ORDER BY rank
                LIMIT 3
            """, (fts_query,))
            results = c.fetchall()
            if not results:
                return ""
            parts = ["[知识库上下文]"]
            for content, filename, chunk_idx in results:
                preview = content[:300].replace('\n', ' ')
                parts.append(f"📄 {filename}#{chunk_idx}: {preview}")
            return "\n".join(parts)
        except Exception as e:
            logger.debug("RAG prefetch failed: %s", e)
            return ""

    def queue_prefetch(self, query: str, *, session_id: str = "") -> None:
        pass  # Synchronous prefetch is fast enough

    def sync_turn(self, user_content: str, assistant_content: str, *,
                  session_id: str = "", messages=None) -> None:
        pass  # RAG doesn't store turns — that's the session DB's job

    def get_tool_schemas(self) -> List[Dict[str, Any]]:
        return [
            {
                "name": "memory_search",
                "description": (
                    "Search the knowledge base (documents indexed via FTS5). "
                    "Returns matching text chunks from indexed files. "
                    "Use this to find information from uploaded documents."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {
                            "type": "string",
                            "description": "Search query — keywords or phrases to find in the knowledge base.",
                        },
                        "limit": {
                            "type": "integer",
                            "default": 5,
                            "description": "Max number of chunks to return (1-10).",
                        },
                    },
                    "required": ["query"],
                },
            },
            {
                "name": "memory_ingest",
                "description": (
                    "Index a file into the knowledge base for future retrieval. "
                    "Supported: txt, md, py, js, json, yaml, html, csv, sh, sql. "
                    "The file will be chunked and indexed with FTS5."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Absolute path to the file to index.",
                        },
                    },
                    "required": ["file_path"],
                },
            },
        ]

    def handle_tool_call(self, tool_name: str, args: Dict[str, Any], **kwargs) -> str:
        if tool_name == "memory_search":
            return self._handle_search(args)
        elif tool_name == "memory_ingest":
            return self._handle_ingest(args)
        return json.dumps({"error": f"Unknown tool: {tool_name}"})

    def _handle_search(self, args: Dict[str, Any]) -> str:
        """Unified ``memory_search`` (logical-unify Slice 2/3).

        Returns curated L1 notes (from the fabric index) plus L4 reference hits
        (RAG documents + any external KB) via the federated L4 hook. This is the
        single path the ``memory_search`` tool uses; bugs like Bug 1 (notes
        invisible to search) are structurally impossible because notes and
        reference stores all flow through ``memory_fabric``.
        """
        query = args.get("query", "").strip()
        limit = min(max(args.get("limit", 5), 1), 10)
        if not query:
            return json.dumps({"error": "query is required"})
        results = []

        # L1 curated notes — single canonical index (fixes Bug 1).
        try:
            from agent.memory_fabric import (
                L1_NOTE,
                recall as fabric_recall,
            )

            for hit in fabric_recall(query, layer=L1_NOTE, limit=limit):
                content = hit.get("content") or ""
                results.append({
                    "doc_id": None,
                    "chunk_id": None,
                    "filename": hit.get("pointer"),
                    "chunk_index": 0,
                    "content": content,
                    "preview": content[:200].replace("\n", " "),
                })
        except Exception:
            logger.debug("memory_fabric L1 recall failed (notes skipped): %s", query)

        # L4 reference — live federation (RAG + external KBs) via the injected
        # hook (memory_manager.search_all). When no hook is wired (legacy
        # single-provider setup), fall back to this provider's own RAG chunks.
        from agent.memory_fabric import get_l4_federation_hook

        hook = get_l4_federation_hook()
        if hook is not None:
            try:
                for hit in hook(query, limit):
                    content = hit.get("content") or ""
                    results.append({
                        "doc_id": None,
                        "chunk_id": None,
                        "filename": hit.get("source") or hit.get("pointer"),
                        "chunk_index": 0,
                        "content": content,
                        "preview": content[:200].replace("\n", " "),
                    })
            except Exception:
                logger.debug("L4 federation failed (skipped): %s", query)
        else:
            # Fallback: this provider's own RAG document store.
            for hit in self.search(query, limit):
                content = hit.get("content") or ""
                results.append({
                    "doc_id": hit.get("doc_id"),
                    "chunk_id": None,
                    "filename": hit.get("filename"),
                    "chunk_index": hit.get("chunk_index", 0),
                    "content": content,
                    "preview": hit.get(
                        "preview", content[:200].replace("\n", " ")
                    ),
                })

        return json.dumps({"results": results, "count": len(results)}, ensure_ascii=False)

    def _handle_ingest(self, args: Dict[str, Any]) -> str:
        file_path = args.get("file_path", "").strip()
        if not file_path or not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"})
        return self.ingest_file(file_path)

    def _store_embedding(self, c: sqlite3.Cursor, chunk_id: int, content: str) -> None:
        """Generate and store embedding for a chunk (fail-open).

        Called during ingest. If embedding API is unavailable, silently skips —
        FTS5 search still works. If sqlite-vec is not loaded, also skips.
        """
        if _VEC_BACKEND != "sqlite-vec" or not _vec_available:
            return
        try:
            from agent.hybrid_retriever import _get_embedding, _vector_to_blob
            vec = _get_embedding(content)
            if vec:
                c.execute(
                    "INSERT OR REPLACE INTO chunks_vec (chunk_id, embedding) VALUES (?, ?)",
                    (chunk_id, _vector_to_blob(vec))
                )
        except Exception as e:
            logger.debug("Embedding storage skipped for chunk %d: %s", chunk_id, e)

    def ingest_file(self, file_path: str) -> str:
        """Index a file into the RAG database."""
        if not self._initialized:
            _init_db(_get_rag_db())
            self._db_path = _get_rag_db()
            self._initialized = True
        try:
            text = _extract_text(file_path)
        except RuntimeError as exc:
            return json.dumps({"error": str(exc)}, ensure_ascii=False)
        if not text.strip():
            return json.dumps({"error": f"No text extracted from {file_path}"})
        p = Path(file_path)
        chunks = _chunk_text(text, file_ext=p.suffix)
        conn = _get_conn(str(self._db_path))
        c = conn.cursor()
        # Check if document already exists
        c.execute("SELECT id FROM documents WHERE path = ?", (str(p.resolve()),))
        existing = c.fetchone()
        if existing:
            doc_id = existing[0]
            # Delete old chunks
            c.execute("DELETE FROM chunks WHERE doc_id = ?", (doc_id,))
            c.execute("DELETE FROM chunks_fts WHERE rowid IN (SELECT id FROM chunks WHERE doc_id = ?)", (doc_id,))
        else:
            c.execute(
                "INSERT INTO documents (path, filename, file_type, file_size, ingested_at, chunk_count) VALUES (?, ?, ?, ?, ?, ?)",
                (str(p.resolve()), p.name, p.suffix, p.stat().st_size, datetime.now().isoformat(), len(chunks))
            )
            doc_id = c.lastrowid
        # Insert chunks
        for idx, chunk in enumerate(chunks):
            c.execute(
                "INSERT INTO chunks (doc_id, chunk_index, content, char_count) VALUES (?, ?, ?, ?)",
                (doc_id, idx, chunk, len(chunk))
            )
            chunk_id = c.lastrowid
            c.execute("INSERT INTO chunks_fts (rowid, content) VALUES (?, ?)", (chunk_id, chunk))
            self._store_embedding(c, chunk_id, chunk)
        # Update chunk count
        c.execute("UPDATE documents SET chunk_count = ? WHERE id = ?", (len(chunks), doc_id))
        conn.commit()
        logger.info("Ingested %s: %d chunks", p.name, len(chunks))
        self._link_related_documents(doc_id, chunks[0] if chunks else "")
        return json.dumps({
            "status": "ok",
            "filename": p.name,
            "chunks": len(chunks),
            "doc_id": doc_id,
        }, ensure_ascii=False)

    def ingest_content(self, filename: str, content: str, file_type: str = "") -> str:
        """Index text content directly (for file uploads)."""
        if not self._initialized:
            _init_db(_get_rag_db())
            self._db_path = _get_rag_db()
            self._initialized = True
        if not content.strip():
            return json.dumps({"error": f"No text content to index"})
        ext = '.' + file_type.lstrip('.') if file_type else ''
        chunks = _chunk_text(content, file_ext=ext)
        if not chunks:
            return json.dumps({"error": "No chunks generated"})
        conn = _get_conn(str(self._db_path))
        c = conn.cursor()
        # Use filename as path identifier
        virtual_path = f"upload://{filename}"
        c.execute("SELECT id FROM documents WHERE path = ?", (virtual_path,))
        existing = c.fetchone()
        if existing:
            doc_id = existing[0]
            c.execute("DELETE FROM chunks WHERE doc_id = ?", (doc_id,))
            c.execute("DELETE FROM chunks_fts WHERE rowid IN (SELECT id FROM chunks WHERE doc_id = ?)", (doc_id,))
        else:
            c.execute(
                "INSERT INTO documents (path, filename, file_type, file_size, ingested_at, chunk_count) VALUES (?, ?, ?, ?, ?, ?)",
                (virtual_path, filename, file_type or Path(filename).suffix, len(content), datetime.now().isoformat(), len(chunks))
            )
            doc_id = c.lastrowid
        for idx, chunk in enumerate(chunks):
            c.execute(
                "INSERT INTO chunks (doc_id, chunk_index, content, char_count) VALUES (?, ?, ?, ?)",
                (doc_id, idx, chunk, len(chunk))
            )
            chunk_id = c.lastrowid
            c.execute("INSERT INTO chunks_fts (rowid, content) VALUES (?, ?)", (chunk_id, chunk))
            self._store_embedding(c, chunk_id, chunk)
        c.execute("UPDATE documents SET chunk_count = ? WHERE id = ?", (len(chunks), doc_id))
        conn.commit()
        logger.info("Ingested content %s: %d chunks", filename, len(chunks))
        self._link_related_documents(doc_id, chunks[0] if chunks else "")
        return json.dumps({
            "status": "ok",
            "filename": filename,
            "chunks": len(chunks),
            "doc_id": doc_id,
        }, ensure_ascii=False)

    def _link_related_documents(self, doc_id: int, sample_text: str) -> None:
        """Detect document-to-document relationships via FTS5 overlap and write DAG edges."""
        if not sample_text or len(sample_text) < 20:
            return
        try:
            import sqlite3 as _sqlite3
            from agent.evolution_manager import get_self_model_db, _get_conn as _get_evo_conn
            # Build FTS5 query using CJK trigram sliding window (same as prefetch)
            conn = _get_conn(str(self._db_path))
            c = conn.cursor()
            safe = re.sub(r'[^\w\u4e00-\u9fff\s]', ' ', sample_text[:200]).strip()
            terms = []
            for part in safe.split():
                cjk_chars = re.findall(r'[\u4e00-\u9fff]+', part)
                ascii_part = re.sub(r'[\u4e00-\u9fff]', '', part)
                if len(ascii_part) >= 3:
                    terms.append(ascii_part)
                for seg in cjk_chars:
                    if len(seg) >= 3:
                        for i in range(len(seg) - 2):
                            terms.append(seg[i:i+3])
            terms = terms[:8]
            if not terms:
                return
            fts_query = " OR ".join(f'"{t}"' for t in terms)
            if not fts_query:
                return
            c.execute("""
                SELECT DISTINCT documents.id, documents.filename
                FROM chunks_fts
                JOIN chunks ON chunks.id = chunks_fts.rowid
                JOIN documents ON documents.id = chunks.doc_id
                WHERE chunks_fts MATCH ? AND documents.id != ?
                LIMIT 5
            """, (fts_query, doc_id))
            related = c.fetchall()
            conn.close()
            if not related:
                return
            # Write document->document edges to Evolution DAG
            evo_conn = _get_evo_conn(str(get_self_model_db()))
            ec = evo_conn.cursor()
            from datetime import datetime as _dt
            ts = _dt.now().isoformat()
            for r_id, r_name in related:
                ec.execute(
                    "INSERT INTO relations (source_type, source_id, target_type, target_id, rel_type, weight, timestamp) "
                    "VALUES (?, ?, ?, ?, ?, ?, ?)",
                    ('document', doc_id, 'document', r_id, 'related', 0.5, ts),
                )
            evo_conn.commit()
            evo_conn.close()
            logger.info("Linked doc %d to %d related documents", doc_id, len(related))
        except Exception as e:
            logger.debug("_link_related_documents failed: %s", e)

    def list_documents(self) -> List[Dict[str, Any]]:
        """List all indexed documents."""
        if not self._initialized:
            return []
        conn = _get_conn(str(self._db_path))
        c = conn.cursor()
        c.execute("SELECT id, filename, file_type, file_size, ingested_at, chunk_count FROM documents ORDER BY ingested_at DESC")
        return [
            {"id": row[0], "filename": row[1], "file_type": row[2],
             "file_size": row[3], "ingested_at": row[4], "chunk_count": row[5]}
            for row in c.fetchall()
        ]

    def search(self, query: str, limit: int = 5) -> List[Dict[str, Any]]:
        """Search the knowledge base and return matching chunks with metadata.
        
        Primary path: FTS5 full-text search (always available).
        Optional path: if vector backend is enabled and embeddings exist,
        performs vector KNN search and merges results (vector hits first).
        Fail-open: vector search errors are logged and swallowed.
        """
        if not self._initialized or not query.strip():
            return []
        try:
            conn = _get_conn(str(self._db_path))
            c = conn.cursor()
            safe_query = re.sub(r'[^\w\u4e00-\u9fff\s]', ' ', query).strip()
            if not safe_query:
                return []
            # Reuse CJK trigram logic from prefetch
            terms = []
            for word in safe_query.split():
                cjk_chars = [ch for ch in word if '\u4e00' <= ch <= '\u9fff']
                ascii_part = ''.join(ch for ch in word if ch not in cjk_chars)
                if ascii_part and len(ascii_part) >= 3:
                    terms.append(ascii_part)
                if len(cjk_chars) >= 3:
                    for i in range(len(cjk_chars) - 2):
                        terms.append(cjk_chars[i] + cjk_chars[i+1] + cjk_chars[i+2])
                elif len(cjk_chars) > 0:
                    # <3 CJK chars: try combining with adjacent terms
                    terms.append(''.join(cjk_chars))
            if not terms:
                return []
            fts_query = " OR ".join(f'"{t}"' for t in terms[:8])
            c.execute("""
                SELECT chunks.content, documents.filename, chunks.chunk_index,
                       documents.id, documents.file_type, chunks.char_count
                FROM chunks_fts
                JOIN chunks ON chunks.id = chunks_fts.rowid
                JOIN documents ON documents.id = chunks.doc_id
                WHERE chunks_fts MATCH ?
                ORDER BY rank
                LIMIT ?
            """, (fts_query, limit))
            fts_results = [
                {"content": row[0], "filename": row[1], "chunk_index": row[2],
                 "doc_id": row[3], "file_type": row[4], "char_count": row[5],
                 "preview": row[0][:300].replace('\n', ' ')}
                for row in c.fetchall()
            ]
            # Vector search (A-1 optional): if enabled, merge KNN results
            if _VEC_BACKEND == "sqlite-vec" and _vec_available:
                vec_results = self._vector_search(query, limit)
                if vec_results:
                    # Deduplicate: vector hits not already in FTS results
                    fts_ids = {r["doc_id"] for r in fts_results}
                    merged = [r for r in vec_results if r["doc_id"] not in fts_ids]
                    # Vector hits first (higher confidence), then FTS
                    return merged[:limit] + fts_results[:limit - len(merged)]
            return fts_results
        except Exception as e:
            logger.debug("RAG search failed: %s", e)
            return []

    def _vector_search(self, query: str, limit: int = 5) -> List[Dict[str, Any]]:
        """Vector KNN search using sqlite-vec.

        Requires embeddings to have been stored during ingest.
        Reuses hybrid_retriever._get_embedding() for query embedding generation.
        Fail-open: any error → return empty list, FTS5 results remain.
        """
        try:
            conn = _get_conn(str(self._db_path))
            c = conn.cursor()
            # Check if any embeddings exist
            c.execute("SELECT COUNT(*) FROM chunks_vec")
            if c.fetchone()[0] == 0:
                return []
            # Generate query embedding via hybrid_retriever's provider resolution
            from agent.hybrid_retriever import _get_embedding, _vector_to_blob
            query_vec = _get_embedding(query)
            if not query_vec:
                logger.debug("RAG vector search: embedding unavailable, falling back to FTS5")
                return []
            # KNN search via sqlite-vec
            query_blob = _vector_to_blob(query_vec)
            c.execute(
                "SELECT chunk_id, distance FROM chunks_vec "
                "WHERE embedding MATCH ? AND k = ? ORDER BY distance",
                (query_blob, limit)
            )
            rows = c.fetchall()
            if not rows:
                return []
            # Fetch chunk content for matched IDs
            results = []
            for chunk_id, distance in rows:
                c.execute(
                    "SELECT c.id, c.doc_id, c.chunk_index, c.content, c.char_count, "
                    "d.filename FROM chunks c JOIN documents d ON c.doc_id = d.id "
                    "WHERE c.id = ?",
                    (chunk_id,)
                )
                row = c.fetchone()
                if row:
                    results.append({
                        "chunk_id": row[0],
                        "doc_id": row[1],
                        "chunk_index": row[2],
                        "content": row[3],
                        "char_count": row[4],
                        "filename": row[5],
                        "score": 1.0 - distance,  # distance is 0 (identical) to 2 (opposite)
                        "source": "vector",
                    })
            logger.debug("RAG vector search: %d results for query (len=%d)", len(results), len(query))
            return results
        except Exception as e:
            logger.debug("Vector search failed (fail-open to FTS5): %s", e)
            return []

    def get_document_chunks(self, doc_id: int) -> List[Dict[str, Any]]:
        """Get all chunks of a document for preview."""
        if not self._initialized:
            return []
        conn = _get_conn(str(self._db_path))
        c = conn.cursor()
        c.execute(
            "SELECT id, chunk_index, content, char_count FROM chunks WHERE doc_id = ? ORDER BY chunk_index",
            (doc_id,)
        )
        return [
            {"id": row[0], "chunk_index": row[1], "content": row[2], "char_count": row[3]}
            for row in c.fetchall()
        ]

    def delete_document(self, doc_id: int) -> bool:
        """Delete a document and its chunks."""
        if not self._initialized:
            return False
        conn = _get_conn(str(self._db_path))
        c = conn.cursor()
        c.execute("SELECT id FROM chunks WHERE doc_id = ?", (doc_id,))
        chunk_ids = [row[0] for row in c.fetchall()]
        for cid in chunk_ids:
            c.execute("DELETE FROM chunks_fts WHERE rowid = ?", (cid,))
            # Also clean up vector embeddings if table exists
            if _VEC_BACKEND == "sqlite-vec" and _vec_available:
                try:
                    c.execute("DELETE FROM chunks_vec WHERE chunk_id = ?", (cid,))
                except Exception as e:
                    logger.debug("chunks_vec cleanup skipped: %s", e)
        c.execute("DELETE FROM chunks WHERE doc_id = ?", (doc_id,))
        c.execute("DELETE FROM documents WHERE id = ?", (doc_id,))
        conn.commit()
        return c.rowcount > 0

    def get_document_stats(self) -> List[Dict[str, Any]]:
        """Return per-document usage stats by querying Evolution DB relations."""
        if not self._initialized:
            return []
        try:
            # Query RAG DB for document list
            conn = _get_conn(str(self._db_path))
            c = conn.cursor()
            c.execute("SELECT id, filename, file_type, chunk_count FROM documents ORDER BY id")
            docs = [
                {"doc_id": row[0], "filename": row[1], "file_type": row[2],
                 "chunk_count": row[3], "query_count": 0}
                for row in c.fetchall()
            ]
            if not docs:
                return []
            # Query Evolution DB for relation counts
            from agent.evolution_manager import get_self_model_db
            evo_conn = sqlite3.connect(str(get_self_model_db()))
            ec = evo_conn.cursor()
            for doc in docs:
                ec.execute(
                    "SELECT COUNT(*) FROM relations WHERE target_type='document' AND target_id=? AND rel_type='queried'",
                    (doc["doc_id"],)
                )
                doc["query_count"] = ec.fetchone()[0]
            evo_conn.close()
            return docs
        except Exception as e:
            logger.debug("get_document_stats failed: %s", e)
            return []

    def shutdown(self) -> None:
        pass
